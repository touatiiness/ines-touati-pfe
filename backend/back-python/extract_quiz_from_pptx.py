from pptx import Presentation
from pathlib import Path
import json
import re

def extract_raw_text_from_docx(docx_path):
    """
    Extrait TOUT le texte brut d'un fichier Word
    DÉTECTE le surlignage jaune pour identifier les bonnes réponses
    """
    try:
        from docx import Document
        from docx.enum.text import WD_COLOR_INDEX
        doc = Document(str(docx_path))

        full_text = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Vérifier si le paragraphe contient du texte surligné en jaune
                has_yellow_highlight = False
                for run in para.runs:
                    # Vérifier plusieurs types de surlignage jaune
                    if run.font.highlight_color in [WD_COLOR_INDEX.YELLOW, 7]:  # 7 = YELLOW
                        has_yellow_highlight = True
                        break
                    # Vérifier aussi la couleur de fond (shading)
                    if hasattr(run, '_element') and run._element.rPr is not None:
                        shd = run._element.rPr.find('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}shd')
                        if shd is not None:
                            fill = shd.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fill')
                            # Détecter les codes couleur jaune (FFFF00, FFFF99, etc.)
                            if fill and fill.upper() in ['FFFF00', 'FFFF99', 'FFFFCC', 'FFFFE0', 'FFEB3B', 'FDD835']:
                                has_yellow_highlight = True
                                break

                # Marquer la ligne surlignée avec ✅
                if has_yellow_highlight:
                    # Détecter si c'est une option (A, B, C, D)
                    if re.match(r'^[A-Fa-f][\)\.]', text):
                        full_text.append(f"{text} ✅ [BONNE RÉPONSE]")
                    else:
                        full_text.append(f"{text} ✅")
                else:
                    full_text.append(text)

        # Joindre tout le texte
        raw_content = "\n".join(full_text)

        return {
            'success': True,
            'raw_text': raw_content,
            'file': docx_path.name if hasattr(docx_path, 'name') else str(docx_path)
        }
    except Exception as e:
        return {
            'success': False,
            'error': str(e),
            'raw_text': '',
            'file': docx_path.name if hasattr(docx_path, 'name') else str(docx_path)
        }


def extract_raw_text_from_pptx(pptx_path):
    """
    Extrait TOUT le texte brut d'un fichier PowerPoint
    DÉTECTE le surlignage jaune pour identifier les bonnes réponses
    """
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))

        full_text = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            full_text.append(f"\n--- Slide {slide_num} ---")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    shape_text = shape.text.strip()

                    # Vérifier si la forme a un fond jaune (fill)
                    has_yellow_bg = False
                    if hasattr(shape, 'fill') and shape.fill.type is not None:
                        try:
                            if hasattr(shape.fill, 'fore_color'):
                                rgb = shape.fill.fore_color.rgb
                                # Détecter les nuances de jaune (RGB proche de 255, 255, 0)
                                if rgb and rgb[0] > 200 and rgb[1] > 200 and rgb[2] < 100:
                                    has_yellow_bg = True
                        except:
                            pass

                    # Vérifier chaque ligne du texte pour détecter le surlignage au niveau du run
                    lines = shape_text.split('\n')
                    processed_lines = []

                    for line in lines:
                        line_stripped = line.strip()
                        if not line_stripped:
                            continue

                        line_has_highlight = has_yellow_bg

                        # Vérifier le surlignage au niveau des runs de texte
                        if hasattr(shape, 'text_frame'):
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.text.strip() and run.text.strip() in line_stripped:
                                        # Vérifier la couleur de fond du run
                                        try:
                                            if hasattr(run.font, 'fill') and run.font.fill.type is not None:
                                                if hasattr(run.font.fill, 'fore_color'):
                                                    rgb = run.font.fill.fore_color.rgb
                                                    if rgb and rgb[0] > 200 and rgb[1] > 200 and rgb[2] < 100:
                                                        line_has_highlight = True
                                        except:
                                            pass

                                        # Vérifier aussi le highlight_color
                                        try:
                                            if hasattr(run.font, 'highlight_color') and run.font.highlight_color:
                                                # Jaune = 7 dans l'énumération MSO_COLOR_TYPE
                                                if run.font.highlight_color == 7:
                                                    line_has_highlight = True
                                        except:
                                            pass

                        # Marquer les réponses surlignées
                        if line_has_highlight and re.match(r'^[A-Fa-f][\)\.]', line_stripped):
                            processed_lines.append(f"{line_stripped} ✅ [BONNE RÉPONSE]")
                        elif line_has_highlight:
                            processed_lines.append(f"{line_stripped} ✅")
                        else:
                            processed_lines.append(line_stripped)

                    full_text.extend(processed_lines)

            # Ajouter les notes si disponibles
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    full_text.append(f"[Notes: {notes}]")

        raw_content = "\n".join(full_text)

        return {
            'success': True,
            'raw_text': raw_content,
            'file': pptx_path.name if hasattr(pptx_path, 'name') else str(pptx_path)
        }
    except Exception as e:
        return {
            'success': False,
            'error': str(e),
            'raw_text': '',
            'file': pptx_path.name if hasattr(pptx_path, 'name') else str(pptx_path)
        }


def extract_quiz_from_docx(docx_path):
    """
    Extrait les questions QCM d'un fichier Word (.docx)

    Format attendu :
    - Texte avec questions numérotées
    - Options A), B), C), D)
    - Peut contenir la réponse

    Returns:
        dict avec 'success', 'total_questions', 'questions', 'formatted_content'
    """
    try:
        from docx import Document
        doc = Document(str(docx_path))
        questions = []
        current_question = None
        question_num = 0

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # Détecter une nouvelle question (commence par un chiffre ou "Question")
            if re.match(r'^\d+[\.\)]\s*|^Question\s*\d+', text, re.IGNORECASE):
                if current_question and current_question['options']:
                    questions.append(current_question)

                question_num += 1
                # Nettoyer le texte de la question
                question_text = re.sub(r'^\d+[\.\)]\s*|^Question\s*\d+[\.\:\s]*', '', text, flags=re.IGNORECASE)

                current_question = {
                    'numero': question_num,
                    'question': question_text,
                    'options': [],
                    'reponse_correcte': None
                }

            # Détecter une option (A), B), etc.)
            elif re.match(r'^[A-Fa-f][\)\.]', text):
                if current_question:
                    current_question['options'].append(text)

            # Chercher réponse correcte (Réponse: A ou Correct: B)
            elif re.search(r'(?:r[ée]ponse|correct|answer)[:=\s]*([A-Fa-f])', text, re.IGNORECASE):
                if current_question:
                    match = re.search(r'(?:r[ée]ponse|correct|answer)[:=\s]*([A-Fa-f])', text, re.IGNORECASE)
                    current_question['reponse_correcte'] = match.group(1).upper()

        # Ajouter la dernière question
        if current_question and current_question['options']:
            questions.append(current_question)

        # Formater le contenu pour le prompt système
        formatted_content = ""
        if questions:
            formatted_content = "\n\n--- QUESTIONS DU QUIZ EXISTANT ---\n"
            for q in questions:
                formatted_content += f"\nQuestion {q['numero']}: {q['question']}\n"
                for option in q['options']:
                    formatted_content += f"{option}\n"
                if q['reponse_correcte']:
                    formatted_content += f"[RÉPONSE CORRECTE: {q['reponse_correcte']}]\n"
            formatted_content += "\n--- FIN DU QUIZ ---\n"

        return {
            'success': True,
            'file': docx_path.name if hasattr(docx_path, 'name') else str(docx_path),
            'total_questions': len(questions),
            'questions': questions,
            'formatted_content': formatted_content
        }

    except Exception as e:
        return {
            'success': False,
            'error': str(e),
            'file': docx_path.name if hasattr(docx_path, 'name') else str(docx_path),
            'total_questions': 0,
            'questions': [],
            'formatted_content': ''
        }


def extract_quiz_from_pptx(file_path):
    """
    Extrait les questions QCM d'un fichier PowerPoint (.pptx) ou Word (.docx)

    Détecte automatiquement le type de fichier et utilise la bonne méthode.

    Format attendu dans les slides PowerPoint:
    - Titre = Question
    - Points de texte = Options (A, B, C, D, etc.)
    - Réponse correcte marquée ou dans les notes

    Format attendu dans Word:
    - Questions numérotées
    - Options A), B), C), D)

    Returns:
        dict avec 'success', 'total_questions', 'questions', 'formatted_content'
    """
    # Détecter le type de fichier
    if str(file_path).endswith('.docx'):
        return extract_quiz_from_docx(file_path)

    # Sinon, traiter comme PowerPoint
    try:
        pptx_path = file_path
        prs = Presentation(str(pptx_path))
        questions = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            question_data = {
                'numero': slide_num,
                'question': '',
                'options': [],
                'reponse_correcte': None,
                'raw_text': []
            }

            # Extraire tout le texte de la slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    question_data['raw_text'].append(text)

                    # Si c'est un titre, c'est probablement la question
                    if shape.is_placeholder:
                        placeholders = slide.placeholders
                        if shape in placeholders and shape.placeholder_format.type == 1:  # Title
                            question_data['question'] = text
                        else:
                            # Texte du corps - probablement des options
                            # Chercher les options A), B), C), etc.
                            lines = text.split('\n')
                            for line in lines:
                                line = line.strip()
                                # Détecter les options avec pattern A), B), a), b), etc.
                                if re.match(r'^[A-Fa-f][\)\.]', line):
                                    question_data['options'].append(line)

                    # Même si pas placeholder, chercher des options dans tout le texte
                    if not shape.is_placeholder or (shape.is_placeholder and shape.placeholder_format.type != 1):
                        lines = text.split('\n')
                        for line in lines:
                            line = line.strip()
                            if re.match(r'^[A-Fa-f][\)\.]', line):
                                if line not in question_data['options']:
                                    question_data['options'].append(line)

            # Extraire les notes du présentateur (peuvent contenir la réponse)
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    # Chercher la réponse dans les notes
                    # Pattern: "Réponse: A" ou "Correct: B" ou juste "A"
                    match = re.search(r'(?:r[ée]ponse|correct|answer)[:=\s]*([A-Fa-f])', notes_text, re.IGNORECASE)
                    if match:
                        question_data['reponse_correcte'] = match.group(1).upper()

            # Si on n'a pas trouvé de question claire, utiliser le premier texte
            if not question_data['question'] and question_data['raw_text']:
                question_data['question'] = question_data['raw_text'][0]

            # Ajouter la question seulement si elle a du contenu
            if question_data['question'] and question_data['options']:
                questions.append(question_data)

        # Formater le contenu pour le prompt système
        formatted_content = ""
        if questions:
            formatted_content = "\n\n--- QUESTIONS DU QUIZ EXISTANT ---\n"
            for q in questions:
                formatted_content += f"\nQuestion {q['numero']}: {q['question']}\n"
                for option in q['options']:
                    formatted_content += f"{option}\n"
                if q['reponse_correcte']:
                    formatted_content += f"[RÉPONSE CORRECTE: {q['reponse_correcte']}]\n"
            formatted_content += "\n--- FIN DU QUIZ ---\n"

        return {
            'success': True,
            'file': pptx_path.name if hasattr(pptx_path, 'name') else str(pptx_path),
            'total_questions': len(questions),
            'questions': questions,
            'formatted_content': formatted_content
        }

    except Exception as e:
        return {
            'success': False,
            'error': str(e),
            'file': pptx_path.name if hasattr(pptx_path, 'name') else str(pptx_path),
            'total_questions': 0,
            'questions': [],
            'formatted_content': ''
        }

def test_extract_quiz(course_num, part_num):
    """
    Teste l'extraction de quiz pour un cours/partie spécifique
    """
    base_path = Path("./Support_Cours_Préparation")
    part_folder = f"{course_num}.{part_num}"
    quizz_folder = base_path / str(course_num) / part_folder / "Quizz"

    if not quizz_folder.exists():
        print(f"❌ Dossier Quizz introuvable: {quizz_folder}")
        return None

    # Trouver le premier fichier .pptx ou .docx
    pptx_files = list(quizz_folder.glob("*.pptx")) + list(quizz_folder.glob("*.docx"))

    if not pptx_files:
        print(f"❌ Aucun fichier de quiz trouvé dans: {quizz_folder}")
        return None

    quiz_file = pptx_files[0]
    print(f"\n{'='*60}")
    print(f"Extraction du quiz: {quiz_file.name}")
    print(f"Cours {course_num}.{part_num}")
    print(f"{'='*60}\n")

    result = extract_quiz_from_pptx(quiz_file)

    if result['success']:
        print(f"✅ Extraction réussie!")
        print(f"   Nombre de questions: {result['total_questions']}\n")

        for q in result['questions']:
            print(f"Question {q['numero']}: {q['question'][:80]}...")
            print(f"   Options: {len(q['options'])} choix")
            if q['reponse_correcte']:
                print(f"   Réponse correcte: {q['reponse_correcte']}")
            print()
    else:
        print(f"❌ Erreur: {result['error']}")

    return result

if __name__ == "__main__":
    # Test avec le cours 1, partie 2 (qui a un quiz visible dans votre screenshot)
    test_extract_quiz(1, 2)
