from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel
import os
import sys
import io
from typing import List, Optional
from pathlib import Path
from openai import OpenAI
from dotenv import load_dotenv

# Charger les variables d'environnement depuis .env
load_dotenv()

# Import des routes de recommandation
from recommendation_routes import router as recommendation_router

# Configuration UTF-8 pour Windows
if sys.platform == "win32":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

app = FastAPI(title="Learning Platform API")

# Inclure les routes de recommandation
app.include_router(recommendation_router)

# Configuration CORS pour permettre les requêtes depuis Angular
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:4200"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configuration OpenAI avec la nouvelle API
client = OpenAI(
    api_key=os.getenv("OPENAI_API_KEY")
)

# Chemin vers les fichiers de cours
COURS_BASE_PATH = Path("./Support_Cours_Préparation")


# Modèles Pydantic
class ChatMessage(BaseModel):
    role: str
    content: str


class ChatRequest(BaseModel):
    course_number: int
    part_number: int
    messages: List[ChatMessage]


class ChatResponse(BaseModel):
    response: str
    conversation: List[ChatMessage]


# Route de test
@app.get("/")
async def root():
    return {"message": "Learning Platform API is running"}


# Route pour l'AI Assistant
@app.post("/api/chat", response_model=ChatResponse)
async def chat_with_ai(request: ChatRequest):
    """
    Endpoint pour discuter avec l'assistant IA pour une partie spécifique d'un cours
    """
    try:
        system_message = {
            "role": "system",
            "content": f"""You are an intelligent teaching assistant for Course {request.course_number}, Part {request.part_number}.
            Your role is to help students understand concepts, answer their questions and guide them in their learning.
            Be patient, clear and encourage the student to think for themselves. Always respond in French."""
        }

        messages = [system_message]
        for msg in request.messages:
            messages.append({"role": msg.role, "content": msg.content})

        # Utiliser la nouvelle API OpenAI
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=messages,
            temperature=0.7,
            max_tokens=500
        )

        ai_response = response.choices[0].message.content

        updated_messages = request.messages.copy()
        updated_messages.append(ChatMessage(role="assistant", content=ai_response))

        return ChatResponse(
            response=ai_response,
            conversation=updated_messages
        )

    except Exception as e:
        print(f"Error with OpenAI: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error communicating with AI: {str(e)}")


# Route pour le Quiz avec AI
@app.post("/api/quiz-chat", response_model=ChatResponse)
async def quiz_chat_with_ai(request: ChatRequest):
    """
    Endpoint pour un quiz interactif géré par l'IA
    """
    try:
        from extract_quiz_from_pptx import extract_raw_text_from_docx, extract_raw_text_from_pptx
        import logging

        # Configuration du logging
        logging.basicConfig(level=logging.INFO)
        logger = logging.getLogger(__name__)

        # Charger les questions du quiz
        part_folder = f"{request.course_number}.{request.part_number}"

        # Essayer avec "Quizz" et "Quiz" (certains dossiers ont des noms différents)
        quizz_folder = COURS_BASE_PATH / str(request.course_number) / part_folder / "Quizz"
        if not quizz_folder.exists():
            quizz_folder = COURS_BASE_PATH / str(request.course_number) / part_folder / "Quiz"

        quiz_raw_content = ""
        has_quiz_file = False

        # Vérifier si des fichiers quiz existent
        if quizz_folder.exists():
            # Chercher fichiers PPTX ET DOCX
            pptx_files = list(quizz_folder.glob("*.pptx"))
            docx_files = list(quizz_folder.glob("*.docx"))

            # Filtrer les fichiers temporaires Word (~$...)
            docx_files = [f for f in docx_files if not f.name.startswith('~$')]

            quiz_files = pptx_files + docx_files

            logger.info(f"📂 Recherche quiz pour Cours {request.course_number}.{request.part_number}")
            logger.info(f"   Dossier: {quizz_folder}")
            logger.info(f"   Fichiers PPTX trouvés: {len(pptx_files)}")
            logger.info(f"   Fichiers DOCX trouvés: {len(docx_files)}")

            # Extraire le TEXTE BRUT pour que l'IA le comprenne
            if quiz_files:
                quiz_file = quiz_files[0]
                logger.info(f"   📄 Lecture du fichier: {quiz_file.name}")

                try:
                    # Extraire texte brut selon le type de fichier
                    if quiz_file.suffix == '.docx':
                        result = extract_raw_text_from_docx(quiz_file)
                    else:  # .pptx
                        result = extract_raw_text_from_pptx(quiz_file)

                    if result['success'] and result['raw_text'].strip():
                        has_quiz_file = True
                        quiz_raw_content = result['raw_text']

                        logger.info(f"   ✅ Texte extrait: {len(quiz_raw_content)} caractères")
                        logger.info(f"      Aperçu: {quiz_raw_content[:100]}...")
                    else:
                        error_msg = result.get('error', 'Fichier vide')
                        logger.warning(f"   ⚠️ Extraction échouée: {error_msg}")
                        logger.warning(f"   → Basculement vers génération de questions par IA")

                except Exception as e:
                    logger.error(f"   ❌ Erreur lors de la lecture: {str(e)}")
                    logger.warning(f"   → Basculement vers génération de questions par IA")
            else:
                logger.info(f"   ℹ️ Aucun fichier PPTX/DOCX trouvé")
                logger.info(f"   → Génération de questions par IA")
        else:
            logger.info(f"📂 Dossier Quiz(z) inexistant: {quizz_folder}")
            logger.info(f"   → Génération de questions par IA")

        # Déterminer le nom du cours pour générer des questions
        course_name = f"Cours {request.course_number} - Partie {request.part_number}"

        # Message système pour l'AI
        if has_quiz_file and quiz_raw_content:
            # Cas 1: Quiz existant - L'IA doit LIRE et COMPRENDRE le contenu
            system_content = f"""Tu es un assistant pédagogique EXPERT qui fait passer un QUIZ pour le {course_name}.

📄 CONTENU DU FICHIER QUIZ :
{quiz_raw_content}

🎯 TON RÔLE :

ÉTAPE 1 - PRÉPARATION (avant de poser la première question):
- LIS et ANALYSE tout le contenu ci-dessus TRÈS ATTENTIVEMENT
- IDENTIFIE toutes les questions et leurs options (A, B, C, D, etc.)
- Pour CHAQUE question, DÉTERMINE la bonne réponse en suivant cette méthode:

  **MÉTHODE 1 - Chercher le marqueur (priorité):**
  - Cherche ✅ [BONNE RÉPONSE] ou "Bonne réponse:" ou "Correct:" dans le texte
  - Si trouvé → Note cette réponse comme la bonne

  **MÉTHODE 2 - Analyser le code (si pas de marqueur OU pour vérifier):**
  - EXÉCUTE mentalement le code ligne par ligne
  - TRACE les valeurs des variables
  - ÉVALUE les conditions (true/false)
  - DÉTERMINE ce qui sera affiché/le résultat
  - Compare avec les options A, B, C, D
  - Note la réponse correcte basée sur ton analyse

  ⚠️ RÈGLE IMPORTANTE: Si tu ne vois PAS de marqueur ✅, utilise TOUJOURS ton analyse du code pour trouver la bonne réponse

- MÉMORISE le nombre total de questions ET la bonne réponse de chaque question

ÉTAPE 2 - PRÉSENTATION:
- Présente-toi chaleureusement
- Indique le nombre total de questions
- Pose la PREMIÈRE question

ÉTAPE 3 - ÉVALUATION (pour CHAQUE réponse de l'étudiant):

🔍 FORMAT DE RÉPONSE:
- L'étudiant répond avec UNE LETTRE: "A", "a", "B", "b", "C", "c", ou "D", "d"
- Compare UNIQUEMENT LA LETTRE (ignore majuscule/minuscule)
- Exemples valides: "c", "C", "reponse c", "je choisis C"
- Extrait la lettre de la réponse de l'étudiant

⚠️ RÈGLES D'ÉVALUATION STRICTES:

**ÉTAPE 1 - RAPPELLE-TOI la bonne réponse que tu as déterminée en ÉTAPE 1:**
Tu as déjà analysé cette question AVANT de la poser. Tu connais déjà la bonne réponse.

**ÉTAPE 2 - ÉVALUATION DE LA RÉPONSE DE L'ÉTUDIANT:**

1. EXTRAIT la lettre de la réponse de l'étudiant (A, B, C ou D)
   - Exemples: "A", "a", "je pense A", "réponse c" → Extrait la lettre

2. COMPARE avec la bonne réponse que TU AS DÉTERMINÉE
   - IGNORE majuscule/minuscule lors de la comparaison
   - Exemple: "A" = "a" → C'est pareil

3. SI LES LETTRES CORRESPONDENT:
   ✅ Dis "✅ Correct!" ou "🎉 Bravo!"
   ✅ Explique POURQUOI en exécutant le code étape par étape:
      - "Car a = 5"
      - "La condition a > 3 est vraie (5 > 3)"
      - "Donc on entre dans le if"
      - "On exécute printf(\"A\\n\")"
      - "La sortie est A"
   ✅ Incrémente le score (+1)

4. SI LES LETTRES NE CORRESPONDENT PAS:
   ❌ Dis "❌ Incorrect."
   ❌ Indique la VRAIE bonne réponse (celle que tu as déterminée)
   ❌ Explique POURQUOI en exécutant le code ligne par ligne
   ❌ Montre où l'étudiant s'est trompé
   ❌ Le score reste inchangé (0 point)

🚫 RÈGLES ABSOLUES - TRÈS IMPORTANT:
- NE JAMAIS dire "Correct" si la lettre est différente !
- NE JAMAIS dire "Incorrect" puis donner une mauvaise explication
- TOUJOURS être COHÉRENT entre ton verdict (Correct/Incorrect) et ton explication
- Si tu dis "La bonne réponse est D", ton explication doit prouver que c'est D (pas A !)

💡 EXEMPLES DE BONNES RÉPONSES:

**EXEMPLE 1 - Étudiant a raison:**
Question: `int a = 5; if (a > 3) printf("A"); else printf("B");`
Analyse mentale: a=5 → 5>3=vrai → if → affiche "A" → Bonne réponse = A
Étudiant répond: "A"
TOI: "✅ Correct! En effet, a vaut 5, la condition 5 > 3 est vraie, donc on entre dans le if et on affiche A."

**EXEMPLE 2 - Étudiant a tort:**
Question: `int a = 5; if (a > 3) printf("A"); else printf("B");`
Analyse mentale: a=5 → 5>3=vrai → if → affiche "A" → Bonne réponse = A
Étudiant répond: "B"
TOI: "❌ Incorrect. La bonne réponse est A. Voici pourquoi : a vaut 5, la condition 5 > 3 est vraie, donc on entre dans le bloc if (et non le else). On exécute printf(\"A\"), donc la sortie est A, pas B."

⚠️ ANTI-EXEMPLE (NE JAMAIS FAIRE ÇA):
Étudiant: "A"
TOI: "❌ Incorrect. La bonne réponse est D: Rien ne s'affiche. Explication: a vaut 5, 5 > 3 est vrai, donc on affiche A."
↑ INCOHÉRENT! Si on affiche A, la réponse est A, pas D!

ÉTAPE 4 - PROGRESSION:

⚠️ VALIDATION DE LA RÉPONSE:
- SI la réponse contient clairement une lettre (A, B, C, ou D): ÉVALUE et PASSE IMMÉDIATEMENT À LA QUESTION SUIVANTE
- SI la réponse est ambiguë (comme "12", "AR", "je sais pas", texte sans lettre, etc.):
  Dis: "⚠️ Je n'ai pas compris votre réponse. Veuillez répondre avec une seule lettre : A, B, C ou D."
  RESTE sur la même question et attends une nouvelle réponse

🔥 RÈGLE ABSOLUE:
- Après avoir évalué UNE réponse (correcte OU incorrecte), passe IMMÉDIATEMENT à la question suivante
- NE RÉPÈTE JAMAIS la même évaluation si l'étudiant renvoie la même réponse
- Si l'étudiant répond "A" puis envoie encore "A", ne redis PAS "❌ Incorrect. La bonne réponse est C..."
- À la place, ignore la répétition et pose directement la question suivante
- NE DEMANDE JAMAIS "Êtes-vous prêt?" - pose directement la question suivante
- Indique TOUJOURS la progression (Question X/Y)
- Pose UNE SEULE question à la fois

ÉTAPE 5 - RÉSULTAT FINAL:
- Après la dernière question, donne le score final : "📊 Score final: X/Y (XX%)"
- Donne des encouragements adaptés au score
- SI LE SCORE < 80%:
  ⚠️ Dis: "Score insuffisant. Le seuil de maîtrise est de 80%."
  📺 RECOMMANDE IMMÉDIATEMENT la vidéo explicative:
  "Je vous recommande de regarder la vidéo explicative de ce cours pour mieux comprendre les concepts."
  "Cliquez ici pour regarder la vidéo du Cours {request.course_number}.{request.part_number}"
- SI LE SCORE ≥ 80%:
  🎉 Félicite chaleureusement l'étudiant pour sa maîtrise du sujet

⚠️ RÈGLES CRITIQUES - IMPÉRATIF:

1. **COHÉRENCE ABSOLUE**: Ton verdict ET ton explication doivent correspondre
   - Si tu dis "La bonne réponse est A", ton explication doit prouver que c'est A
   - JAMAIS: "La bonne réponse est D" + "donc on affiche A" (INCOHÉRENT!)

2. **ANALYSE RIGOUREUSE DU CODE C**:
   - Exécute le code mentalement ligne par ligne
   - Trace toutes les variables
   - Évalue toutes les conditions (true/false)
   - Détermine le résultat AVANT de répondre

3. **PRIORITÉS**:
   - Si marqueur ✅ visible → Respecte-le
   - Sinon → Analyse le code avec tes connaissances (tu es un expert en C)

4. **PÉDAGOGIE**: Explique TOUJOURS en détaillant l'exécution du code

5. **FRANÇAIS**: Réponds TOUJOURS en français

6. **NE JAMAIS**: Féliciter si la réponse est fausse, ou critiquer si elle est correcte

🔥 COMMENCE LE QUIZ MAINTENANT!"""
        else:
            # Cas 2: Aucun quiz - l'IA doit générer ses propres questions
            system_content = f"""Tu es un assistant pédagogique intelligent qui fait passer un QUIZ pour le {course_name}.

ATTENTION: Aucun quiz prédéfini n'est disponible pour ce cours.

Ton rôle:
1. Présente-toi et explique qu'il y a un quiz à faire sur le sujet du {course_name}
2. GÉNÈRE TOI-MÊME 5-7 questions pertinentes sur ce sujet (format QCM avec 4 options A, B, C, D)
3. Pose les questions UNE PAR UNE (ne pose jamais toutes les questions en même temps)
4. Attends la réponse de l'étudiant
5. Évalue la réponse (correcte/incorrecte)
6. Si incorrect, explique pourquoi et donne la bonne réponse avec des explications pédagogiques
7. Si correct, félicite l'étudiant et explique brièvement pourquoi c'est correct
8. Passe à la question suivante
9. À la fin, donne le score total et des encouragements

IMPORTANT:
- GÉNÈRE des questions pertinentes et éducatives sur le sujet
- Pose SEULEMENT UNE question à la fois
- Attends la réponse avant de poser la suivante
- Garde un compte du score
- Sois encourageant et pédagogique
- Réponds TOUJOURS en français
- Utilise les émojis pour rendre le quiz plus engageant
- Assure-toi que tes questions couvrent les concepts clés du sujet"""

        system_message = {
            "role": "system",
            "content": system_content
        }

        messages = [system_message]
        for msg in request.messages:
            messages.append({"role": msg.role, "content": msg.content})

        # Utiliser la nouvelle API OpenAI
        # GPT-4 est beaucoup plus fiable pour l'analyse de code et la cohérence
        response = client.chat.completions.create(
            model="gpt-4",
            messages=messages,
            temperature=0.3,  # Plus bas pour plus de précision
            max_tokens=800
        )

        ai_response = response.choices[0].message.content

        updated_messages = request.messages.copy()
        updated_messages.append(ChatMessage(role="assistant", content=ai_response))

        return ChatResponse(
            response=ai_response,
            conversation=updated_messages
        )

    except Exception as e:
        print(f"Error with Quiz AI: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error communicating with AI: {str(e)}")


# Route pour obtenir les fichiers de cours (adaptée à la structure avec Cours/ et Quizz/)
@app.get("/api/courses/{course_number}/parts/{part_number}/files")
async def get_course_files(course_number: int, part_number: int, file_type: str = "Cours"):
    """
    Retourne la liste des fichiers disponibles pour une partie d'un cours
    file_type peut être "Cours" ou "Quizz"

    Structure : Support_Cours_Préparation/1/1.1/Cours/fichier.pptx
    """
    try:
        # Construire le chemin : 1/1.1/Cours ou 1/1.1/Quizz
        part_folder = f"{course_number}.{part_number}"
        folder_path = COURS_BASE_PATH / str(course_number) / part_folder / file_type

        if not folder_path.exists():
            raise HTTPException(status_code=404, detail=f"Dossier {file_type} introuvable pour ce cours")

        # Lister les fichiers
        files = []
        for file_path in folder_path.iterdir():
            if file_path.is_file():
                files.append({
                    "name": file_path.name,
                    "type": file_path.suffix,
                    "size": file_path.stat().st_size,
                    "path": str(file_path.relative_to(COURS_BASE_PATH)),
                    "category": file_type
                })

        return {"files": files}

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erreur: {str(e)}")


# Route pour lister TOUS les fichiers (Cours + Quizz)
@app.get("/api/courses/{course_number}/parts/{part_number}/all-files")
async def get_all_course_files(course_number: int, part_number: int):
    """
    Retourne tous les fichiers (Cours ET Quizz) pour une partie
    """
    try:
        part_folder = f"{course_number}.{part_number}"
        base_folder = COURS_BASE_PATH / str(course_number) / part_folder

        if not base_folder.exists():
            raise HTTPException(status_code=404, detail="Partie du cours introuvable")

        all_files = []

        # Parcourir Cours/
        cours_folder = base_folder / "Cours"
        if cours_folder.exists():
            for file_path in cours_folder.iterdir():
                if file_path.is_file():
                    all_files.append({
                        "name": file_path.name,
                        "type": file_path.suffix,
                        "size": file_path.stat().st_size,
                        "path": str(file_path.relative_to(COURS_BASE_PATH)),
                        "category": "Cours"
                    })

        # Parcourir Quizz/
        quizz_folder = base_folder / "Quizz"
        if quizz_folder.exists():
            for file_path in quizz_folder.iterdir():
                if file_path.is_file():
                    all_files.append({
                        "name": file_path.name,
                        "type": file_path.suffix,
                        "size": file_path.stat().st_size,
                        "path": str(file_path.relative_to(COURS_BASE_PATH)),
                        "category": "Quizz"
                    })

        # Parcourir Video/ (avec ou sans accent)
        video_folder = base_folder / "Video"
        if not video_folder.exists():
            video_folder = base_folder / "Vidéo"  # Essayer avec accent

        print(f"   Vérification dossier Video: {video_folder}")
        print(f"   Existe? {video_folder.exists()}")

        if video_folder.exists():
            print(f"   ✅ Dossier Video trouvé!")
            for file_path in video_folder.iterdir():
                if file_path.is_file():
                    print(f"      📹 Fichier vidéo: {file_path.name} (type: {file_path.suffix})")
                    all_files.append({
                        "name": file_path.name,
                        "type": file_path.suffix,
                        "size": file_path.stat().st_size,
                        "path": str(file_path.relative_to(COURS_BASE_PATH)),
                        "category": "Video"
                    })
        else:
            print(f"   ❌ Dossier Video/Vidéo NON trouvé")

        print(f"   Total fichiers trouvés: {len(all_files)}")
        return {"files": all_files}

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erreur: {str(e)}")


# Route pour télécharger un fichier
@app.get("/api/courses/{course_number}/parts/{part_number}/download/{category}/{file_name}")
async def download_course_file(course_number: int, part_number: int, category: str, file_name: str):
    """
    Télécharge un fichier de cours
    category: "Cours" ou "Quizz"
    """
    try:
        part_folder = f"{course_number}.{part_number}"
        file_path = COURS_BASE_PATH / str(course_number) / part_folder / category / file_name

        if not file_path.exists():
            raise HTTPException(status_code=404, detail="Fichier introuvable")

        return FileResponse(
            path=str(file_path),
            filename=file_name,
            media_type="application/octet-stream"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erreur: {str(e)}")


# Route pour visualiser un fichier
@app.get("/api/courses/{course_number}/parts/{part_number}/view/{category}/{file_name}")
async def view_course_file(course_number: int, part_number: int, category: str, file_name: str):
    """
    Retourne un fichier pour visualisation
    category: "Cours", "Quizz" ou "Video"
    """
    try:
        part_folder = f"{course_number}.{part_number}"
        file_path = COURS_BASE_PATH / str(course_number) / part_folder / category / file_name

        # Si le fichier n'existe pas et category="Video", essayer avec "Vidéo"
        if not file_path.exists() and category == "Video":
            file_path = COURS_BASE_PATH / str(course_number) / part_folder / "Vidéo" / file_name

        if not file_path.exists():
            print(f"❌ Fichier introuvable: {file_path}")
            raise HTTPException(status_code=404, detail="Fichier introuvable")

        print(f"✅ Fichier trouvé: {file_path}")

        mime_types = {
            ".pdf": "application/pdf",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".txt": "text/plain",
            ".jpg": "image/jpeg",
            ".png": "image/png",
            ".mp4": "video/mp4",
            ".avi": "video/x-msvideo",
            ".mkv": "video/x-matroska",
            ".mov": "video/quicktime",
            ".webm": "video/webm"
        }

        file_extension = file_path.suffix.lower()
        media_type = mime_types.get(file_extension, "application/octet-stream")

        print(f"📹 Type MIME: {media_type}")

        # Ajouter le header Content-Disposition pour forcer l'affichage inline
        headers = {}
        if file_extension == ".pdf" or file_extension in [".mp4", ".avi", ".mkv", ".mov", ".webm"]:
            headers["Content-Disposition"] = f'inline; filename="{file_name}"'

        return FileResponse(
            path=str(file_path),
            media_type=media_type,
            filename=file_name,
            headers=headers
        )

    except Exception as e:
        print(f"❌ Erreur: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Erreur: {str(e)}")


# Route pour obtenir les questions de quiz
@app.get("/api/courses/{course_number}/parts/{part_number}/quiz")
async def get_quiz_questions(course_number: int, part_number: int):
    """
    Retourne les questions de quiz pour une partie spécifique d'un cours
    """
    try:
        from pptx import Presentation
        import re

        part_folder = f"{course_number}.{part_number}"
        quizz_folder = COURS_BASE_PATH / str(course_number) / part_folder / "Quizz"

        if not quizz_folder.exists():
            raise HTTPException(status_code=404, detail="Dossier Quizz introuvable")

        # Trouver le premier fichier PowerPoint
        pptx_files = list(quizz_folder.glob("*.pptx"))
        if not pptx_files:
            raise HTTPException(status_code=404, detail="Aucun fichier de quiz trouvé")

        quiz_file = pptx_files[0]
        prs = Presentation(str(quiz_file))
        questions = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            question_text = ""
            options = []
            correct_answer = None

            # Extraire le texte de la slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()

                    # Première forme avec texte = question
                    if not question_text and not re.match(r'^[A-Fa-f][\)\.]', text):
                        question_text = text
                    else:
                        # Extraire les options
                        lines = text.split('\n')
                        for line in lines:
                            line = line.strip()
                            if re.match(r'^[A-Fa-f][\)\.]', line):
                                options.append(line)

            # Extraire la réponse des notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                match = re.search(r'(?:r[ée]ponse|correct|answer)[:=\s]*([A-Fa-f])', notes_text, re.IGNORECASE)
                if match:
                    correct_answer = match.group(1).upper()

            if question_text and options:
                questions.append({
                    'id': slide_num,
                    'question': question_text,
                    'options': options,
                    'correct_answer': correct_answer
                })

        return {
            'quiz_file': quiz_file.name,
            'total_questions': len(questions),
            'questions': questions
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erreur: {str(e)}")


# Route pour évaluer une réponse de quiz
@app.post("/api/courses/{course_number}/parts/{part_number}/quiz/check")
async def check_quiz_answer(course_number: int, part_number: int, question_id: int, user_answer: str):
    """
    Vérifie si la réponse de l'utilisateur est correcte
    """
    try:
        # Récupérer les questions du quiz
        quiz_data = await get_quiz_questions(course_number, part_number)

        # Trouver la question correspondante
        question = next((q for q in quiz_data['questions'] if q['id'] == question_id), None)

        if not question:
            raise HTTPException(status_code=404, detail="Question introuvable")

        correct_answer = question.get('correct_answer')
        user_answer_normalized = user_answer.upper().strip()

        is_correct = correct_answer and user_answer_normalized == correct_answer

        return {
            'is_correct': is_correct,
            'correct_answer': correct_answer,
            'user_answer': user_answer_normalized
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erreur: {str(e)}")


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8001)
